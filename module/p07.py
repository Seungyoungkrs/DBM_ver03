import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import numpy as np
import datetime
import warnings
from pyNastran.op2.op2 import OP2, OP2Writer
from sklearn.metrics import r2_score, mean_squared_error, mean_absolute_percentage_error

# warnings
warnings.simplefilter("ignore")

class DBM_ResultGenerator:
    def __init__(self, analyser):
        """
        Initialize the DBM_ResultGenerator class with a DBM_Analyser instance.
        
        Parameters:
        - analyser: An instance of DBM_Analyser containing the analysis data.
        """
        self.analyser = analyser
        self.freq = analyser.freq
        self.ha = analyser.ha
        self.LC = analyser.LC
        self.ref_sensor_loc = analyser.ref_sensor_loc
        self.ref_sensor_ids = analyser.ref_sensor_ids
        self.export_ppt = analyser.export_ppt
        self.ElemNode_ref_df = analyser.ElemNode_ref_df
        self.Y = analyser.Y
        self.Yref_all = analyser.Yref_all
        self.B = analyser.B
        
        self.all_references = []
        self.all_predictions = []
        self.all_metadata = pd.DataFrame(columns=['LC', 'HA', 'Location', 'Ref', 'ElemID'])

    def cal_statistic(self, LC=None, ha=None, check_elem=None, make_graph=False):
        if LC is None:
            LC = self.LC
        if ha is None:
            ha = self.ha
        if check_elem is None:
            check_elem = self.ref_sensor_ids.to_list()
            
        data1 = []
        nFreq = len(self.freq)
        
        # PPT 프레젠테이션 객체 생성
        if self.export_ppt:
            prs = Presentation()
        
        ref_list = self.ref_sensor_ids.to_list()
        ref_index = [ref_list.index(str(elem)) for elem in check_elem if str(elem) in ref_list]
        
        for rLC in LC:
            for rHA in ha:
                nHA = ha.index(rHA)
                nLC = LC.index(rLC)                
                nHALC = (nHA + 1) + (len(ha) * nLC)
                
                for i in ref_index:
                    # 예측 값(Y_prd)과 실제 값(Y_ref) 추출
                    Y_prd = self.Y.iloc[(nHALC - 1) * nFreq:nHALC * nFreq, i]
                    Y_ref = self.Yref_all.iloc[(nHALC - 1) * nFreq:nHALC * nFreq, i]
                    Y_ref_val = Y_ref.values
                    Y_prd_val = Y_prd.values

                    # 예측 값과 실제 값을 통합 리스트에 추가
                    self.all_predictions.extend(Y_prd.values)
                    self.all_references.extend(Y_ref.values)

                    # 메타데이터 추가
                    location = self.ref_sensor_loc[i]
                    ElemID= self.ElemNode_ref_df.iloc[i, 0]
                    new_metadata = pd.DataFrame({'LC': [rLC], 'HA': [rHA], 'Location': [location], 'Ref': [i + 1], 'ElemID': [ElemID]})
                    self.all_metadata = pd.concat([self.all_metadata, new_metadata], ignore_index=True)

                    # 통계 지표 계산
                    Range = np.max(Y_ref_val) - np.min(Y_ref_val)
                    corrXY = np.corrcoef(Y_prd_val, Y_ref_val)[0, 1]
                    MSE = mean_squared_error(Y_ref_val, Y_prd_val)
                    RMSE = np.sqrt(MSE)                    
                    MAPE = np.mean(np.abs((Y_ref_val - Y_prd_val) / Y_ref_val)) * 100         
                    
                    # Range가 0인 경우가 있음 (ref 요소가 yz평면에 있을 경우 xx응력은 0이라서 발생함)
                    if Range == 0:
                        RMSE_ratio = 0  
                    else:
                        RMSE_ratio = np.sqrt(np.mean(np.square((Y_ref_val - Y_prd_val) / Range)))
                    r2 = r2_score(Y_ref, Y_prd)                    
                    
                    # 좌표 계산
                    posX = (self.ElemNode_ref_df.iloc[i, 5] + self.ElemNode_ref_df.iloc[i, 8] + self.ElemNode_ref_df.iloc[i, 11] + self.ElemNode_ref_df.iloc[i, 14]) / 4
                    posY = (self.ElemNode_ref_df.iloc[i, 6] + self.ElemNode_ref_df.iloc[i, 9] + self.ElemNode_ref_df.iloc[i, 12] + self.ElemNode_ref_df.iloc[i, 15]) / 4
                    posZ = (self.ElemNode_ref_df.iloc[i, 7] + self.ElemNode_ref_df.iloc[i, 10] + self.ElemNode_ref_df.iloc[i, 13] + self.ElemNode_ref_df.iloc[i, 16]) / 4
                    
                    # 데이터 추가
                    data = [rLC, rHA, "All", i + 1, corrXY, RMSE, MAPE, RMSE_ratio, r2]                    
                    data1.append(data)
                    
                    if make_graph == True:
                        plt.figure()
                        plt.plot(self.freq, Y_ref_val, label='Reference')
                        plt.plot(self.freq, Y_prd_val, label='Predict')
                        plt.text(0.03, 0.03, f"R2 = {r2:.3f}\nRMSE = {RMSE:.3f}\nPOS:({posX:.0f},{posY:.0f},{posZ:.0f})\nStation: {location}", fontsize=10, transform=plt.gca().transAxes)
                        plt.title(f"Ref. {i + 1} (Elem. {self.ref_sensor_ids[i]}) ({rLC}/{rHA}°)", pad=10, size=20)
                        plt.legend()
                        plt.grid()
                        plt.xlabel("Frequency (rad/sec)", size=14)
                        plt.ylabel("σxx (MPa)", size=14)

                    #PPT 저장 및 그래프 생성 (미구현)
                    if self.export_ppt:                                            
                        plt.figure()
                        plt.plot(self.freq, Y_ref_val, label='Reference')
                        plt.plot(self.freq, Y_prd_val, label='Predict')
                        plt.text(0.03, 0.03, f"R2 = {r2:.3f}\nRMSE = {RMSE:.3f}\nPOS:({posX:.0f},{posY:.0f},{posZ:.0f})\nStation: {location}", fontsize=10, transform=plt.gca().transAxes)
                        plt.title(f"Ref. {i + 1} (Elem. {self.ref_sensor_ids[i]}) ({rLC}/{rHA}°)", pad=10, size=20)
                        plt.legend()
                        plt.grid()
                        plt.xlabel("Frequency (rad/sec)", size=14)
                        plt.ylabel("σxx (MPa)", size=14)

                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        title = slide.shapes.title
                        title.text = f'HA: {rHA}, LC: {rLC}, Ref: {i + 1}'

                        # 그래프를 이미지로 저장
                        img_path = f'temp_plot_{rHA}_{rLC}_{i + 1}.png'
                        plt.savefig(img_path)
                        plt.close()

                        # PPT 슬라이드에 이미지 삽입
                        slide.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(8), height=Inches(4.5))
                      
        if self.export_ppt:
            prs.save('DBM_Analysis_Presentation.pptx')
            
        # 전체 데이터 셋에 대한 통계 지표 계산
        overall_r2 = r2_score(self.all_references, self.all_predictions)
        overall_rmse = np.sqrt(mean_squared_error(self.all_references, self.all_predictions))
        overall_mape = mean_absolute_percentage_error(self.all_references, self.all_predictions)

        print(f"Overall R²: {overall_r2}")
        print(f"Overall RMSE: {overall_rmse}")
        print(f"Overall MAPE: {overall_mape}")

        # r2, RMSE, MAPE 값만 추출
        r2_values = [item[-1] for item in data1]
        rmse_values = [item[5] for item in data1]
        mape_values = [item[6] for item in data1]

        # 각 값들의 평균 계산
        average_r2 = sum(r2_values) / len(r2_values)
        average_rmse = sum(rmse_values) / len(rmse_values)
        average_mape = sum(mape_values) / len(mape_values)

        print(f"Average r2: {average_r2}")
        print(f"Average RMSE: {average_rmse}")
        print(f"Average MAPE: {average_mape}")
                
#         # 결과 CSV 파일로 저장
#         pd.DataFrame(self.all_references).to_csv('all_references.csv')
#         pd.DataFrame(self.all_predictions).to_csv('all_predictions.csv')

    def get_filtered_data(self, LC=None, HA=None, Location=None, Ref=None, priority='Location'):
        """
        필터링으로 통계데이터를 확인하는 함수 (개발 중)
        """
        filtered_metadata = self.all_metadata.copy()

        # 필터링 조건에 따른 메타데이터 필터링
        if LC is not None:
            if isinstance(LC, list):
                filtered_metadata = filtered_metadata[filtered_metadata['LC'].isin(LC)]
            else:
                filtered_metadata = filtered_metadata[filtered_metadata['LC'] == LC]

        if HA is not None:
            if isinstance(HA, list):
                filtered_metadata = filtered_metadata[filtered_metadata['HA'].isin(HA)]
            else:
                filtered_metadata = filtered_metadata[filtered_metadata['HA'] == HA]

        if Location is not None and (Ref is None or priority == 'Location'):
            if isinstance(Location, list):
                filtered_metadata = filtered_metadata[filtered_metadata['Location'].isin(Location)]
            else:
                filtered_metadata = filtered_metadata[filtered_metadata['Location'] == Location]

        if Ref is not None and (Location is None or priority == 'Ref'):
            if isinstance(Ref, list):
                filtered_metadata = filtered_metadata[filtered_metadata['Ref'].isin(Ref)]
            else:
                filtered_metadata = filtered_metadata[filtered_metadata['Ref'] == Ref]

        filtered_indices = filtered_metadata.index.tolist()
        num_freq = len(self.freq)

        detailed_data = []
        combined_predictions = []
        combined_references = []

        # 필터링된 인덱스를 기반으로 데이터 추출 및 결합
        for i in filtered_indices:
            start_idx = i * num_freq
            end_idx = start_idx + num_freq

            predictions = self.all_predictions[start_idx:end_idx]
            references = self.all_references[start_idx:end_idx]

            combined_predictions.extend(predictions)
            combined_references.extend(references)

            # 세부 데이터 생성
            for j in range(num_freq):
                detailed_data.append({
                    'LC': filtered_metadata.loc[i, 'LC'],
                    'HA': filtered_metadata.loc[i, 'HA'],
                    'Location': filtered_metadata.loc[i, 'Location'],
                    'Ref': filtered_metadata.loc[i, 'Ref'],
                    'Frequency': self.freq[j],
                    'Predictions': predictions[j],
                    'References': references[j]
                })

        # 필터링된 데이터에 대한 통계 재계산
        r2 = r2_score(combined_references, combined_predictions)
        rmse = np.sqrt(mean_squared_error(combined_references, combined_predictions))
        mape = mean_absolute_percentage_error(combined_references, combined_predictions)

        # 결과 데이터프레임 생성
        result_df = pd.DataFrame(detailed_data)
        result_df['R2'] = r2
        result_df['RMSE'] = rmse
        result_df['MAPE'] = mape

        return result_df

    def save_results_to_f06(self, filename='results.f06'):
        """
        f06 파일에 결과 저장
        HullScan에서 읽기 가능 (Global 변환 미구현으로 local xx에서만 Global xx확인 가능)
        """
        now = datetime.datetime.now()
        date_str = now.strftime("%d/%m/%Y")
        time_str = now.strftime("%H:%M:%S")
        created_date_str = now.strftime("%b %d, %Y").upper()
        nastran_version = "MSC Nastran  6/20/23"
        page_num = 1
        el_no = len(self.ElemNode_ref_df)
        freq_no = len(self.freq)

        with open(filename, 'w') as f:
            subcase_id = 1
            for rLC in self.LC:
                nLC = self.LC.index(rLC)  
                for rHA in self.ha:
                    nHA = self.ha.index(rHA)
                    for freq in self.freq:
                        nFreq = self.freq.index(freq)
                        elem_counter = 0  # 요소 카운터 초기화
                        if nLC == 0:
                            start_index = nHA * el_no * freq_no + nFreq   
                        elif nLC == 1:
                            start_index = int(len(self.all_predictions) /2) + nHA * el_no * freq_no + nFreq   
                        if elem_counter % 16 == 0:                                
                            f.write(f"1    MSC.NASTRAN JOB CREATED ON {date_str} AT {time_str}                            {created_date_str}  {nastran_version}   PAGE   {page_num}\n")
                            page_num += 1
                            title = f"PRD{subcase_id}_{rLC}_HA_{rHA}_FREQ_{freq:.2f}"
                            f.write(f"     {title:<116}\n")
                            f.write(f"0                                                                                                            SUBCASE {subcase_id}\n")
                            f.write(" \n")
                            f.write("                         S T R E S S E S   I N   Q U A D R I L A T E R A L   E L E M E N T S   ( Q U A D 4 )\n")
                            f.write("  ELEMENT      FIBER               STRESSES IN ELEMENT COORD SYSTEM             PRINCIPAL STRESSES (ZERO SHEAR)                 \n")
                            f.write("    ID.       DISTANCE           NORMAL-X       NORMAL-Y      SHEAR-XY       ANGLE         MAJOR           MINOR        VON MISES\n")

                        nHALC = (nHA + 1) + (len(self.ha) * nLC)
                        elem_no_list = [self.ref_sensor_ids[ii] for ii in range(el_no)]
                        elem_index = 0
                        result_index = [start_index + ii * freq_no for ii in range(el_no)]                        

                        for j in range(len(result_index)):
                            jj = result_index[j]
                            elem_no = int(elem_no_list[j])
                            pred = self.all_predictions[jj]                        
    #                         if elem_no == 110676:
    #                             print(f"i : {j}, HA : {rHA}, Freq: {nFreq}, elem_no: {elem_no}, Index : {jj}")
    #                             print(pred)
                            normal_x = pred
                            normal_y = 0.0
                            shear_xy = 0.0
                            angle = 0.0
                            major = normal_x
                            minor = normal_x
                            von_mises = abs(normal_x)

                            f.write(f"0 {elem_no:<8} -6.000000E+00  {normal_x:>16.6E} {normal_y:>14.6E} {shear_xy:>14.6E} {angle:>12.4f} {major:>16.6E} {minor:>16.6E} {von_mises:>14.6E}\n")
                            f.write(f" {'':<9}  6.000000E+00  {normal_x:>16.6E} {normal_y:>14.6E} {shear_xy:>14.6E} {angle:>12.4f} {major:>16.6E} {minor:>16.6E} {von_mises:>14.6E}\n")
                            elem_index += 1
                            elem_counter += 1

                        subcase_id += 1
        print(f"Results saved to {filename}")
